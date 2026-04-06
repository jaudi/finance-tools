"""pages/13_Pitch_Deck.py — Pitch Deck Builder
Generates a 9-slide investor pitch deck as a .pptx file.
"""
import io
from datetime import date

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Page config ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Pitch Deck Builder · Finance Tools",
    page_icon="🎯",
    layout="wide",
)

st.markdown("""<style>
    .block-container { padding-top: 1.5rem; }
    h2 { font-size:1.25rem!important; color:#333;
         border-bottom:2px solid #eee; padding-bottom:6px; margin-top:1.5rem; }
    .ev-caption { color:#888; font-size:0.88rem; margin-top:-6px; margin-bottom:1rem; }
</style>""", unsafe_allow_html=True)

from mobile_css import inject_mobile_css
inject_mobile_css()
st.page_link("app.py", label="← All Tools")

st.title("🎯 Pitch Deck Builder")
st.markdown(
    '<p class="ev-caption">9-slide investor pitch deck · Light professional template · '
    'PowerPoint export (.pptx) · Open in PowerPoint or Google Slides</p>',
    unsafe_allow_html=True,
)
st.info(
    "Fill in each tab below, then click **Generate Pitch Deck** at the bottom to download your .pptx file.",
    icon="ℹ️",
)

# ══════════════════════════════════════════════════════════════════════════════
# PPTX CONSTANTS & HELPERS
# ══════════════════════════════════════════════════════════════════════════════
_C_PRIMARY = RGBColor(0x00, 0x3F, 0x88)
_C_ACCENT  = RGBColor(0x00, 0x66, 0xCC)
_C_LIGHT   = RGBColor(0xF7, 0xF9, 0xFC)
_C_BORDER  = RGBColor(0xE8, 0xED, 0xF4)
_C_DARK    = RGBColor(0x22, 0x22, 0x22)
_C_GREY    = RGBColor(0x88, 0x88, 0x88)
_C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
_C_GREEN   = RGBColor(0x1A, 0x93, 0x6F)
_C_RED     = RGBColor(0xC0, 0x39, 0x2B)

_RECT = 1  # MSO_AUTO_SHAPE_TYPE — plain rectangle


def _tb(slide, text, l, t, w, h, size=12, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or _C_DARK


def _rect(slide, l, t, w, h, fill):
    """Add a filled rectangle with no visible border."""
    shape = slide.shapes.add_shape(_RECT, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill  # match border to fill → invisible
    return shape


def _bar(slide):
    """Left blue accent bar used on content slides."""
    _rect(slide, 0, 0, 0.09, 7.5, _C_PRIMARY)


def _header(slide, title, sub=""):
    """Header band with title and optional subtitle."""
    _rect(slide, 0.09, 0, 13.24, 1.35, _C_LIGHT)
    _tb(slide, title, 0.3, 0.13, 12.5, 0.85, size=26, bold=True, color=_C_PRIMARY)
    if sub:
        _tb(slide, sub, 0.3, 0.95, 12.5, 0.35, size=10, italic=True, color=_C_GREY)


def _footer(slide, co):
    """Bottom footer band."""
    _rect(slide, 0.09, 7.2, 13.24, 0.3, _C_LIGHT)
    _tb(slide, f"{co}  ·  Confidential", 0.3, 7.22, 9.0, 0.25, size=8, color=_C_GREY)
    _tb(slide, "financeplots.com", 9.5, 7.22, 3.6, 0.25,
        size=8, color=_C_GREY, align=PP_ALIGN.RIGHT)


def _cell_fmt(tbl, r, c, text, size=11, bold=False, color=None,
              align=PP_ALIGN.CENTER, bg=None):
    """Format a table cell."""
    cell = tbl.cell(r, c)
    cell.text = str(text)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or _C_DARK
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


# ══════════════════════════════════════════════════════════════════════════════
# PPTX BUILDER
# ══════════════════════════════════════════════════════════════════════════════
def build_pptx(d: dict) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    co    = d["company_name"]
    sym   = d["currency_sym"]

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.33, 7.5, _C_WHITE)       # background
    _rect(s, 0, 0, 0.22, 7.5, _C_PRIMARY)      # left bar (wider on cover)
    _rect(s, 0.22, 6.3, 13.11, 1.2, _C_LIGHT)  # bottom strip

    if d["logo_bytes"]:
        try:
            s.shapes.add_picture(
                io.BytesIO(d["logo_bytes"]), Inches(10.8), Inches(0.3), height=Inches(0.9)
            )
        except Exception:
            pass

    _tb(s, co, 0.45, 2.15, 11.5, 1.3, size=44, bold=True, color=_C_PRIMARY)
    _rect(s, 0.45, 3.55, 2.8, 0.06, _C_ACCENT)
    _tb(s, d["tagline"], 0.45, 3.7, 11.5, 0.75, size=20, color=_C_ACCENT)
    _tb(s, f"{d['industry']}  ·  {d['deck_date'].strftime('%B %Y')}",
        0.45, 6.4, 10.5, 0.4, size=11, color=_C_GREY)
    _tb(s, "Strictly Confidential — Not for Distribution",
        0.45, 6.85, 10.5, 0.35, size=9, italic=True, color=_C_GREY)

    # ── SLIDE 2: PROBLEM ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bar(s)
    _header(s, "The Problem", "Pain points your company solves")
    _footer(s, co)
    probs = [(lbl, desc) for lbl, desc in d["problems"] if lbl]
    if probs:
        n  = len(probs)
        bw = (12.5 - 0.15 * (n - 1)) / n
        for i, (lbl, desc) in enumerate(probs):
            lx = 0.3 + i * (bw + 0.15)
            _rect(s, lx, 1.5, bw, 5.55, _C_LIGHT)
            _rect(s, lx, 1.5, bw, 0.5, _C_PRIMARY)
            _tb(s, str(i + 1), lx + 0.12, 1.55, bw - 0.24, 0.4,
                size=16, bold=True, color=_C_WHITE)
            _tb(s, lbl, lx + 0.15, 2.1, bw - 0.3, 0.55, size=14, bold=True, color=_C_PRIMARY)
            _rect(s, lx + 0.15, 2.7, bw - 0.3, 0.04, _C_ACCENT)
            _tb(s, desc, lx + 0.15, 2.82, bw - 0.3, 4.1, size=12, color=_C_DARK)

    # ── SLIDE 3: SOLUTION ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bar(s)
    _header(s, "Our Solution", "How we uniquely solve the problem")
    _footer(s, co)
    _rect(s, 0.3, 1.5, 12.5, 2.05, _C_LIGHT)
    _rect(s, 0.3, 1.5, 0.06, 2.05, _C_ACCENT)
    _tb(s, d["sol_narrative"], 0.55, 1.65, 12.1, 1.75, size=14, color=_C_DARK)
    diffs = [x for x in d["differentiators"] if x]
    if diffs:
        _tb(s, "Key Differentiators", 0.3, 3.68, 12.5, 0.4, size=11, bold=True, color=_C_GREY)
        dw = (12.5 - 0.15 * (len(diffs) - 1)) / len(diffs)
        for i, diff in enumerate(diffs):
            dx = 0.3 + i * (dw + 0.15)
            _rect(s, dx, 4.18, dw, 2.6, _C_LIGHT)
            _rect(s, dx, 4.18, dw, 0.07, _C_ACCENT)
            _tb(s, diff, dx + 0.15, 4.38, dw - 0.3, 2.27, size=13, bold=True, color=_C_PRIMARY)

    # ── SLIDE 4: MARKET ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bar(s)
    _header(s, "Market Opportunity", "TAM · SAM · SOM")
    _footer(s, co)
    mw = (12.5 - 0.3) / 3
    for i, (lbl, val, desc, clr) in enumerate([
        ("TAM", f"{sym}{d['tam_val']:,.1f}B", "Total Addressable Market",      _C_PRIMARY),
        ("SAM", f"{sym}{d['sam_val']:,.1f}B", "Serviceable Addressable Market", _C_ACCENT),
        ("SOM", f"{sym}{d['som_val']:,.0f}M", "Serviceable Obtainable Market",  _C_GREEN),
    ]):
        mx = 0.3 + i * (mw + 0.15)
        _rect(s, mx, 1.5, mw, 4.6, _C_LIGHT)
        _rect(s, mx, 1.5, mw, 0.07, clr)
        _tb(s, lbl,  mx + 0.2, 1.72, mw - 0.4, 0.5,  size=13, bold=True, color=clr)
        _tb(s, val,  mx + 0.2, 2.28, mw - 0.4, 1.05, size=30, bold=True, color=_C_PRIMARY)
        _tb(s, desc, mx + 0.2, 3.4,  mw - 0.4, 0.55, size=10, color=_C_GREY)
    if d["mkt_narrative"]:
        _tb(s, d["mkt_narrative"], 0.3, 6.25, 12.5, 0.8, size=11, color=_C_GREY)

    # ── SLIDE 5: BUSINESS MODEL ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bar(s)
    _header(s, "Business Model", "How we make money")
    _footer(s, co)
    for xi, (title, body) in enumerate([
        ("Revenue Streams", d["rev_streams"]),
        ("Pricing",         d["pricing"]),
    ]):
        lx = 0.3 + xi * 6.55
        _rect(s, lx, 1.5, 6.2, 5.55, _C_LIGHT)
        _tb(s, title, lx + 0.2, 1.65, 5.8, 0.5, size=13, bold=True, color=_C_PRIMARY)
        _rect(s, lx + 0.2, 2.2, 5.8, 0.04, _C_ACCENT)
        for j, line in enumerate([ln for ln in body.split("\n") if ln.strip()][:9]):
            _tb(s, f"•  {line}", lx + 0.2, 2.35 + j * 0.5, 5.8, 0.45, size=11, color=_C_DARK)

    # ── SLIDE 6: TRACTION ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bar(s)
    _header(s, "Traction & Milestones", "Proof this is working")
    _footer(s, co)
    kpis = [(lbl, val) for lbl, val in d["kpis"] if lbl]
    if kpis:
        kw = (12.5 - 0.15 * (len(kpis) - 1)) / len(kpis)
        for i, (kl, kv) in enumerate(kpis):
            kx = 0.3 + i * (kw + 0.15)
            _rect(s, kx, 1.5, kw, 2.1, _C_LIGHT)
            _rect(s, kx, 1.5, kw, 0.07, _C_ACCENT)
            _tb(s, kv, kx + 0.15, 1.7,  kw - 0.3, 0.95, size=28, bold=True, color=_C_PRIMARY)
            _tb(s, kl, kx + 0.15, 2.7,  kw - 0.3, 0.45, size=10, color=_C_GREY)
    _tb(s, "Key Milestones", 0.3, 3.75, 12.5, 0.45, size=12, bold=True, color=_C_PRIMARY)
    _rect(s, 0.3, 4.24, 12.5, 0.04, _C_ACCENT)
    for j, m in enumerate([ln for ln in d["milestones"].split("\n") if ln.strip()][:6]):
        _tb(s, f"✓  {m}", 0.3, 4.38 + j * 0.45, 12.5, 0.42, size=11, color=_C_DARK)

    # ── SLIDE 7: FINANCIALS ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bar(s)
    _header(s, "Financial Projections", f"5-year forecast · {sym}")
    _footer(s, co)
    rv = d["rev_vals"]
    pv = d["prof_vals"]

    tbl = s.shapes.add_table(3, 6, Inches(0.3), Inches(1.6), Inches(12.5), Inches(2.3)).table

    for ci, h in enumerate(["", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]):
        _cell_fmt(tbl, 0, ci, h, bold=True, color=_C_WHITE, bg=_C_PRIMARY)

    for ri, (lbl, vals, bg) in enumerate([
        (f"Revenue ({sym})",    rv, _C_LIGHT),
        (f"Net Profit ({sym})", pv, _C_WHITE),
    ], 1):
        _cell_fmt(tbl, ri, 0, lbl, bold=True, color=_C_PRIMARY, align=PP_ALIGN.LEFT, bg=bg)
        for ci, val in enumerate(vals, 1):
            clr = _C_GREEN if val >= 0 else _C_RED
            _cell_fmt(tbl, ri, ci, f"{sym}{val:,.0f}", color=clr, bg=bg)

    growths = []
    for i in range(1, 5):
        if rv[i - 1] > 0:
            growths.append(f"Y{i}→Y{i+1}: +{(rv[i] - rv[i-1]) / rv[i-1] * 100:.0f}%")
    if growths:
        _tb(s, "YoY Revenue Growth:  " + "  |  ".join(growths),
            0.3, 4.1, 12.5, 0.38, size=10, color=_C_GREY)
    if rv[0] > 0:
        cagr = (rv[4] / rv[0]) ** 0.25 - 1
        _tb(s, f"5-Year Revenue CAGR:  {cagr * 100:.1f}%",
            0.3, 4.6, 6.5, 0.48, size=14, bold=True, color=_C_ACCENT)

    # ── SLIDE 8: TEAM ─────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bar(s)
    _header(s, "The Team", "The people behind the mission")
    _footer(s, co)
    team = d["team"]
    if team:
        n  = min(len(team), 4)
        tw = (12.5 - 0.15 * (n - 1)) / n
        for i, (name, role, bio) in enumerate(team[:4]):
            tx = 0.3 + i * (tw + 0.15)
            _rect(s, tx, 1.5, tw, 5.55, _C_LIGHT)
            _rect(s, tx, 1.5, tw, 0.07, _C_PRIMARY)
            # Initials avatar placeholder
            _rect(s, tx + tw / 2 - 0.5, 1.7, 1.0, 1.0, _C_BORDER)
            initials = "".join(w[0].upper() for w in name.split()[:2])
            _tb(s, initials, tx + tw / 2 - 0.45, 1.88, 0.9, 0.65,
                size=18, bold=True, color=_C_PRIMARY, align=PP_ALIGN.CENTER)
            _tb(s, name, tx + 0.15, 2.85, tw - 0.3, 0.5, size=13, bold=True, color=_C_PRIMARY)
            _tb(s, role, tx + 0.15, 3.38, tw - 0.3, 0.4, size=10, color=_C_ACCENT)
            _rect(s, tx + 0.15, 3.85, tw - 0.3, 0.04, _C_BORDER)
            _tb(s, bio, tx + 0.15, 3.97, tw - 0.3, 2.93, size=10, color=_C_DARK)

    # ── SLIDE 9: THE ASK ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bar(s)
    _header(s, "The Ask", d["funding_type"] + " Round")
    _footer(s, co)

    # Left panel — amount + use of funds
    _tb(s, f"{sym}{d['funding_amt']:,.0f}",
        0.3, 1.55, 6.3, 1.2, size=42, bold=True, color=_C_PRIMARY)
    _tb(s, d["funding_type"] + " Round",
        0.3, 2.78, 6.3, 0.5, size=16, color=_C_ACCENT)
    _rect(s, 0.3, 3.4, 6.2, 0.04, _C_BORDER)
    _tb(s, "Use of Funds", 0.3, 3.55, 6.2, 0.4, size=11, bold=True, color=_C_PRIMARY)
    uof = [(lbl, pct) for lbl, pct in d["use_of_funds"] if lbl]
    if uof:
        uw = (6.2 - 0.15 * (len(uof) - 1)) / len(uof)
        for i, (ul, up) in enumerate(uof):
            ux = 0.3 + i * (uw + 0.15)
            _rect(s, ux, 4.05, uw, 2.1, _C_LIGHT)
            _rect(s, ux, 4.05, max(0.05, uw * up / 100), 0.12, _C_ACCENT)
            _tb(s, f"{up}%", ux + 0.1, 4.28, uw - 0.2, 0.65,
                size=20, bold=True, color=_C_PRIMARY)
            _tb(s, ul, ux + 0.1, 4.98, uw - 0.2, 0.65, size=9, color=_C_GREY)

    # Right panel — what you'll achieve
    _rect(s, 6.75, 1.5, 6.25, 5.6, _C_LIGHT)
    _rect(s, 6.75, 1.5, 6.25, 0.07, _C_GREEN)
    _tb(s, "With this investment we will:", 6.95, 1.65, 5.85, 0.45,
        size=12, bold=True, color=_C_PRIMARY)
    for j, ach in enumerate([ln for ln in d["achieve_text"].split("\n") if ln.strip()][:7]):
        _tb(s, f"→  {ach}", 6.95, 2.25 + j * 0.55, 5.85, 0.5, size=12, color=_C_DARK)

    # ── Export ────────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════════════════
# STREAMLIT FORM
# ══════════════════════════════════════════════════════════════════════════════
tabs = st.tabs([
    "1 · Cover", "2 · Problem", "3 · Solution", "4 · Market",
    "5 · Business Model", "6 · Traction", "7 · Financials", "8 · Team", "9 · The Ask",
])

# ── Tab 1: Cover ──────────────────────────────────────────────────────────────
with tabs[0]:
    st.header("Cover Slide")
    c1, c2 = st.columns(2)
    with c1:
        company_name = st.text_input("Company name *", value="My Company")
        tagline      = st.text_input("Tagline / one-liner *", value="Transforming X with Y")
        industry     = st.text_input("Industry", value="SaaS")
    with c2:
        currency_opt = st.selectbox("Currency (used in financials)", ["USD ($)", "EUR (€)", "GBP (£)"])
        deck_date    = st.date_input("Deck date", value=date.today())
        logo_file    = st.file_uploader("Company logo — PNG / JPG (optional)", type=["png", "jpg", "jpeg"])

currency_sym = currency_opt.split("(")[1].replace(")", "")
logo_bytes   = logo_file.read() if logo_file else None

# ── Tab 2: Problem ────────────────────────────────────────────────────────────
with tabs[1]:
    st.header("The Problem")
    st.caption("Describe up to 3 pain points your company solves.")
    p1c, p2c, p3c = st.columns(3)
    with p1c:
        p1_label = st.text_input("Problem 1 label *", value="Inefficiency")
        p1_desc  = st.text_area("Problem 1 description", value="Manual processes waste time and money.", height=110)
    with p2c:
        p2_label = st.text_input("Problem 2 label *", value="High Cost")
        p2_desc  = st.text_area("Problem 2 description", value="Existing solutions are expensive and hard to implement.", height=110)
    with p3c:
        p3_label = st.text_input("Problem 3 label (optional)", value="")
        p3_desc  = st.text_area("Problem 3 description (optional)", value="", height=110)

problems = [(p1_label, p1_desc), (p2_label, p2_desc), (p3_label, p3_desc)]

# ── Tab 3: Solution ───────────────────────────────────────────────────────────
with tabs[2]:
    st.header("Our Solution")
    sol_narrative = st.text_area(
        "Solution narrative *",
        value="We provide an easy-to-use platform that automates X, reduces costs by Y, and delivers Z.",
        height=130,
    )
    st.caption("Up to 3 key differentiators shown as highlight cards:")
    d1c, d2c, d3c = st.columns(3)
    with d1c:
        diff1 = st.text_input("Differentiator 1", value="10× faster")
    with d2c:
        diff2 = st.text_input("Differentiator 2", value="50% cheaper")
    with d3c:
        diff3 = st.text_input("Differentiator 3 (optional)", value="")

differentiators = [diff1, diff2, diff3]

# ── Tab 4: Market ─────────────────────────────────────────────────────────────
with tabs[3]:
    st.header("Market Opportunity")
    m1c, m2c, m3c = st.columns(3)
    with m1c:
        tam_val = st.number_input(f"TAM ({currency_sym}B — Total Addressable Market)", min_value=0.0, value=10.0, step=0.5)
    with m2c:
        sam_val = st.number_input(f"SAM ({currency_sym}B — Serviceable Addressable)", min_value=0.0, value=2.0, step=0.1)
    with m3c:
        som_val = st.number_input(f"SOM ({currency_sym}M — Serviceable Obtainable)", min_value=0.0, value=50.0, step=5.0)
    mkt_narrative = st.text_area("Market narrative (optional — shown below the boxes)", value="", height=80)

# ── Tab 5: Business Model ─────────────────────────────────────────────────────
with tabs[4]:
    st.header("Business Model")
    bm1c, bm2c = st.columns(2)
    with bm1c:
        rev_streams = st.text_area(
            "Revenue streams (one per line)",
            value="SaaS subscription (monthly / annual)\nProfessional services\nMarketplace transaction fees",
            height=140,
        )
    with bm2c:
        pricing = st.text_area(
            "Pricing tiers (one per line)",
            value="Starter: $49/mo\nGrowth: $199/mo\nEnterprise: custom",
            height=140,
        )

# ── Tab 6: Traction ───────────────────────────────────────────────────────────
with tabs[5]:
    st.header("Traction & Milestones")
    st.caption("Up to 4 headline metrics:")
    k1c, k2c, k3c, k4c = st.columns(4)
    with k1c:
        kpi1_label = st.text_input("KPI 1 label", value="MRR")
        kpi1_val   = st.text_input("KPI 1 value", value="$25K")
    with k2c:
        kpi2_label = st.text_input("KPI 2 label", value="Customers")
        kpi2_val   = st.text_input("KPI 2 value", value="120")
    with k3c:
        kpi3_label = st.text_input("KPI 3 label (optional)", value="MoM Growth")
        kpi3_val   = st.text_input("KPI 3 value (optional)", value="18%")
    with k4c:
        kpi4_label = st.text_input("KPI 4 label (optional)", value="")
        kpi4_val   = st.text_input("KPI 4 value (optional)", value="")
    milestones = st.text_area(
        "Key milestones (one per line)",
        value="Launched MVP — Jan 2024\nFirst 100 paying customers — Apr 2024\nSeed round closed — Jul 2024",
        height=120,
    )

kpis = [
    (kpi1_label, kpi1_val), (kpi2_label, kpi2_val),
    (kpi3_label, kpi3_val), (kpi4_label, kpi4_val),
]

# ── Tab 7: Financials ─────────────────────────────────────────────────────────
with tabs[6]:
    st.header("Financial Projections (5 Years)")
    st.caption(f"Annual figures in {currency_sym}.")
    rev_vals  = []
    prof_vals = []
    yr_cols   = st.columns(5)
    for i, col in enumerate(yr_cols, 1):
        with col:
            st.markdown(f"**Year {i}**")
            rv = st.number_input(
                f"Revenue ({currency_sym})", min_value=0, value=i * 200_000,
                step=10_000, key=f"rv{i}",
            )
            pv = st.number_input(
                f"Net Profit ({currency_sym})", value=int(i * 200_000 * 0.05 * i),
                step=5_000, key=f"pv{i}",
            )
            rev_vals.append(rv)
            prof_vals.append(pv)

# ── Tab 8: Team ───────────────────────────────────────────────────────────────
with tabs[7]:
    st.header("The Team")
    st.caption("Up to 4 members.")
    team_members = []
    row1 = st.columns(2)
    row2 = st.columns(2)
    _defaults = [
        (1, row1[0], "John Smith",  "CEO & Co-Founder",  "10 years in B2B SaaS. Previously VP at Acme Corp."),
        (2, row1[1], "Jane Doe",    "CTO & Co-Founder",  "Ex-Google engineer. 8 years in ML infrastructure."),
        (3, row2[0], "", "", ""),
        (4, row2[1], "", "", ""),
    ]
    for idx, col, def_name, def_role, def_bio in _defaults:
        with col:
            nm = st.text_input(f"Name {idx}", value=def_name, key=f"tn{idx}")
            rl = st.text_input(f"Role {idx}", value=def_role, key=f"tr{idx}")
            bi = st.text_area(f"Bio {idx} (1-2 lines)", value=def_bio, height=80, key=f"tb{idx}")
            if nm:
                team_members.append((nm, rl, bi))

# ── Tab 9: The Ask ────────────────────────────────────────────────────────────
with tabs[8]:
    st.header("The Ask")
    ask1c, ask2c = st.columns(2)
    with ask1c:
        funding_amt  = st.number_input(
            f"Funding amount ({currency_sym})", min_value=0, value=1_500_000, step=50_000
        )
        funding_type = st.selectbox(
            "Round type", ["Pre-seed", "Seed", "Series A", "Series B", "Bridge / Convertible"]
        )
        achieve_text = st.text_area(
            "What you'll achieve with this raise (one per line)",
            value="18-month runway\nReach $1M ARR\nExpand to 3 new markets",
            height=120,
        )
    with ask2c:
        st.caption("Use of funds — label + % each (should total 100):")
        uof_items = []
        _uof_defaults = [
            ("Product & Engineering", 40),
            ("Sales & Marketing",     35),
            ("Operations",            25),
            ("",                       0),
        ]
        for fi, (def_l, def_p) in enumerate(_uof_defaults, 1):
            fc1, fc2 = st.columns([3, 1])
            with fc1:
                ul = st.text_input(f"Label {fi}", value=def_l, key=f"uf{fi}l")
            with fc2:
                up = st.number_input("%", min_value=0, max_value=100, value=def_p, key=f"uf{fi}p")
            uof_items.append((ul, up))

# ══════════════════════════════════════════════════════════════════════════════
# GENERATE BUTTON
# ══════════════════════════════════════════════════════════════════════════════
st.markdown("---")
if st.button("🎯  Generate Pitch Deck (.pptx)", type="primary", use_container_width=True):
    with st.spinner("Building your pitch deck…"):
        pptx_buf = build_pptx({
            "company_name":    company_name,
            "tagline":         tagline,
            "industry":        industry,
            "currency_sym":    currency_sym,
            "deck_date":       deck_date,
            "logo_bytes":      logo_bytes,
            "problems":        problems,
            "sol_narrative":   sol_narrative,
            "differentiators": differentiators,
            "tam_val":         tam_val,
            "sam_val":         sam_val,
            "som_val":         som_val,
            "mkt_narrative":   mkt_narrative,
            "rev_streams":     rev_streams,
            "pricing":         pricing,
            "kpis":            kpis,
            "milestones":      milestones,
            "rev_vals":        rev_vals,
            "prof_vals":       prof_vals,
            "team":            team_members,
            "funding_amt":     funding_amt,
            "funding_type":    funding_type,
            "use_of_funds":    uof_items,
            "achieve_text":    achieve_text,
        })
    st.success("Pitch deck ready — click below to download.")
    st.download_button(
        label="⬇  Download Pitch Deck (.pptx)",
        data=pptx_buf,
        file_name=f"{company_name.replace(' ', '_')}_Pitch_Deck.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )

st.markdown("---")
st.markdown(
    "<small>Built by <b>FinancePlots</b> · "
    "Exported as an editable .pptx — open in PowerPoint or Google Slides · "
    "For informational purposes only</small>",
    unsafe_allow_html=True,
)
